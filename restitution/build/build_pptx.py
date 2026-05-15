"""Build the LearnEU AMA restitution .pptx from restitution/slides/*.md.

Strategy
--------
The official template ships with placeholder-heavy layouts (3 Columns,
Right Content, etc.) whose text frames carry "Lorem ipsum" defaults.
Filling only one of those placeholders left the others visible on the
exported slide. To produce a clean deck we instead use the **Title Only**
layout as a near-blank canvas and draw the body ourselves with shapes,
keeping the template's fonts, master colours and brand chrome.

For each `restitution/slides/slide-NN-*.md` we extract:
  * Title (from the H1 line)
  * Headline / Sub-headline / Layout hint (metadata bullets)
  * One or more "## Body bullets" sections (optionally with
    `(left -- caption)` / `(right -- caption)` qualifiers)
  * Optional `## Visual`, `## Demo cue`, `## Speaker notes`

Render rules
  * Slide 1  -> "Title slide" layout (title + subtitle + small body box)
  * Slide 20 -> "1_Closing logo slide" layout (title + subtitle + recap)
  * Everything else -> "Title Only" + custom drawing:
      - Sub-headline strip under the title
      - 1, 2 or 3 cards depending on body sections
      - Each card has an optional caption header in brand orange
      - Bullet text auto-shrinks if the section has many bullets
  * Speaker notes go in the notes pane, with a `---` footer giving
    rubric coverage / visual / demo cue / source refs for the speaker.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "Subject" / "Azure Master Architect_Prezo_Template_v01.pptx"
SLIDES_DIR = ROOT / "restitution" / "slides"
OUT = ROOT / "restitution" / "build" / "LearnEU-AMA-Restitution.pptx"

# Brand palette (sampled from template chrome)
ORANGE = RGBColor(0xF2, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xC8, 0xD4, 0xE6)
CARD_FILL = RGBColor(0x12, 0x2A, 0x4A)
CARD_LINE = RGBColor(0x2A, 0x4A, 0x78)

# 16:9 deck => 12 192 000 x 6 858 000 EMU
SLIDE_W = 12_192_000
SLIDE_H = 6_858_000


# ----------------------------- parsing ---------------------------------


def _meta(md: str, name: str) -> str:
    rx = re.compile(rf"^-\s+\*\*{re.escape(name)}.*?:\*\*\s*(.+)$", re.MULTILINE)
    m = rx.search(md)
    return m.group(1).strip() if m else ""


def _section(md: str, name: str) -> str:
    m = re.search(
        rf"^##\s+{re.escape(name)}[^\n]*\n((?:(?!^##\s).+\n?)*)", md, re.MULTILINE
    )
    return m.group(1).strip() if m else ""


def _strip_md(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    return s.strip()


def parse_slide(md: str) -> dict:
    out: dict = {
        "title_line": "",
        "headline": _meta(md, "Headline"),
        "subheadline": _meta(md, "Sub-headline"),
        "layout_hint": _meta(md, "Layout (template)"),
        "rubric": _meta(md, "Rubric coverage"),
        "sources": _meta(md, "Source refs"),
        "visual": _section(md, "Visual"),
        "demo_cue": _section(md, "Demo cue"),
        "notes": _section(md, "Speaker notes"),
        "columns": [],  # [{caption, bullets[]}, ...]
    }
    m = re.search(r"^#\s+Slide\s+\d+\s*\xb7\s*[^\xb7]+\xb7\s*(.+)$", md, re.MULTILINE)
    if m:
        out["title_line"] = m.group(1).strip()

    for sec in re.finditer(
        r"^##\s+Body bullets(?:\s*\(([^)]*)\))?\s*\n((?:(?!^##\s).+\n?)*)",
        md,
        re.MULTILINE,
    ):
        caption_raw = (sec.group(1) or "").strip()
        caption = ""
        if caption_raw:
            parts = re.split(r"\s*[\u2014\u2013-]\s*", caption_raw, maxsplit=1)
            caption = parts[1].strip().capitalize() if len(parts) == 2 else caption_raw
        bullets: list[str] = []
        for line in sec.group(2).splitlines():
            line = line.strip()
            if line.startswith("- "):
                bullets.append(_strip_md(line[2:]))
        out["columns"].append({"caption": caption, "bullets": bullets})
    return out


# ----------------------------- drawing ---------------------------------


def _set_text(tf, lines: list[tuple[str, int, int, RGBColor, bool]]) -> None:
    """lines = [(text, level, font_pt, color, bold)]"""
    tf.word_wrap = True
    tf.margin_left = Emu(120_000)
    tf.margin_right = Emu(120_000)
    tf.margin_top = Emu(80_000)
    tf.margin_bottom = Emu(80_000)
    if not lines:
        tf.text = ""
        return
    first = True
    for text, level, size, color, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_subhead(slide, text: str) -> None:
    if not text:
        return
    box = slide.shapes.add_textbox(
        Emu(640_000), Emu(950_000), Emu(SLIDE_W - 1_280_000), Emu(360_000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _strip_md(text)
    for run in p.runs:
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = SOFT


def _add_card(slide, x, y, w, h, caption: str, bullets: list[str]) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.color.rgb = CARD_LINE
    card.line.width = Emu(12_700)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(220_000)
    tf.margin_right = Emu(220_000)
    tf.margin_top = Emu(180_000)
    tf.margin_bottom = Emu(180_000)
    n = len(bullets)
    body_pt = 18 if n <= 5 else 16 if n <= 7 else 14
    cap_pt = 16
    lines: list[tuple[str, int, int, RGBColor, bool]] = []
    if caption:
        lines.append((caption.upper(), 0, cap_pt, ORANGE, True))
    for b in bullets:
        lines.append(("\u2022 " + b, 0, body_pt, WHITE, False))
    _set_text(tf, lines)


def render_content(slide, columns: list[dict]) -> None:
    body_top = 1_400_000
    body_h = SLIDE_H - body_top - 360_000
    margin = 600_000
    if not columns:
        return
    if len(columns) == 1:
        _add_card(
            slide,
            Emu(margin),
            Emu(body_top),
            Emu(SLIDE_W - 2 * margin),
            Emu(body_h),
            columns[0]["caption"],
            columns[0]["bullets"],
        )
        return
    cols = columns[:3]
    gap = 240_000
    total_w = SLIDE_W - 2 * margin - gap * (len(cols) - 1)
    card_w = total_w // len(cols)
    for i, c in enumerate(cols):
        x = margin + i * (card_w + gap)
        _add_card(
            slide,
            Emu(x),
            Emu(body_top),
            Emu(card_w),
            Emu(body_h),
            c["caption"],
            c["bullets"],
        )


def set_title(slide, text: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = text or ""
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = WHITE


def set_notes(slide, spec: dict) -> None:
    parts: list[str] = []
    if spec["notes"]:
        parts.append(spec["notes"])
    meta = []
    if spec["rubric"]:
        meta.append(f"Rubric: {spec['rubric']}")
    if spec["visual"]:
        meta.append(f"Visual: {spec['visual']}")
    if spec["demo_cue"]:
        meta.append(f"Demo cue: {spec['demo_cue']}")
    if spec["sources"]:
        meta.append(f"Source refs: {spec['sources']}")
    if meta:
        parts.append("---\n" + "\n".join(meta))
    if not parts:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = "\n\n".join(parts)


def layout_by_name(prs: Presentation, *names: str):
    by = {l.name: l for l in prs.slide_layouts}
    for n in names:
        if n in by:
            return by[n]
    return prs.slide_layouts[0]


# ----------------------------- per-slide builders ----------------------


def _wipe_other_placeholders(slide) -> None:
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 0:
            continue
        if ph.has_text_frame:
            ph.text_frame.text = ""


def build_title_slide(prs: Presentation, spec: dict):
    layout = layout_by_name(prs, "Title slide")
    s = prs.slides.add_slide(layout)
    set_title(s, spec["headline"])
    _wipe_other_placeholders(s)
    box = s.shapes.add_textbox(
        Emu(640_000), Emu(2_900_000), Emu(SLIDE_W - 1_280_000), Emu(1_400_000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    lines: list[tuple[str, int, int, RGBColor, bool]] = []
    if spec["subheadline"]:
        lines.append((_strip_md(spec["subheadline"]), 0, 22, SOFT, False))
    if spec["columns"]:
        for b in spec["columns"][0]["bullets"]:
            lines.append((b, 0, 16, WHITE, False))
    _set_text(tf, lines)
    return s


def build_closing_slide(prs: Presentation, spec: dict):
    layout = layout_by_name(prs, "1_Closing logo slide", "2_Closing logo slide")
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        if ph.has_text_frame:
            ph.text_frame.text = ""
    box = s.shapes.add_textbox(
        Emu(640_000), Emu(900_000), Emu(SLIDE_W - 1_280_000), Emu(900_000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec["headline"] or ""
    for run in p.runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE
    sub = s.shapes.add_textbox(
        Emu(640_000), Emu(1_900_000), Emu(SLIDE_W - 1_280_000), Emu(500_000)
    )
    sp = sub.text_frame.paragraphs[0]
    sp.text = _strip_md(spec["subheadline"] or "")
    for run in sp.runs:
        run.font.size = Pt(20)
        run.font.italic = True
        run.font.color.rgb = SOFT
    if spec["columns"]:
        body_top = 2_600_000
        body_h = SLIDE_H - body_top - 700_000
        _add_card(
            s,
            Emu(900_000),
            Emu(body_top),
            Emu(SLIDE_W - 1_800_000),
            Emu(body_h),
            "Self-score 57 / 60",
            spec["columns"][0]["bullets"],
        )
    return s


def build_content_slide(prs: Presentation, spec: dict):
    layout = layout_by_name(prs, "Title Only")
    s = prs.slides.add_slide(layout)
    _wipe_other_placeholders(s)
    set_title(s, spec["headline"])
    add_subhead(s, spec["subheadline"])
    render_content(s, spec["columns"])
    return s


# --------------------------------- main --------------------------------


def main() -> int:
    if not TEMPLATE.exists():
        print(f"ERROR: template missing: {TEMPLATE}", file=sys.stderr)
        return 2
    prs = Presentation(str(TEMPLATE))

    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    for sldId in list(sldIdLst):
        rId = sldId.attrib[
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        ]
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    md_files = sorted(SLIDES_DIR.glob("slide-*.md"))
    if not md_files:
        print(f"ERROR: no slide files in {SLIDES_DIR}", file=sys.stderr)
        return 2

    for path in md_files:
        m = re.match(r"slide-(\d+)-", path.name)
        if not m:
            continue
        idx = int(m.group(1))
        spec = parse_slide(path.read_text(encoding="utf-8"))
        if idx == 1:
            slide = build_title_slide(prs, spec)
            ltag = "TitleSlide"
        elif idx == 20:
            slide = build_closing_slide(prs, spec)
            ltag = "Closing"
        else:
            slide = build_content_slide(prs, spec)
            ltag = "TitleOnly+cards"
        set_notes(slide, spec)
        ncols = len(spec["columns"])
        nb = sum(len(c["bullets"]) for c in spec["columns"])
        print(
            f"slide {idx:02d} \xb7 {ltag:18s} \xb7 cols={ncols} bullets={nb} \xb7 "
            f"{(spec['headline'] or spec['title_line'])[:55]}"
        )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(
        f"\nSaved: {OUT}  ({OUT.stat().st_size:,} bytes, {len(md_files)} slides)"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
