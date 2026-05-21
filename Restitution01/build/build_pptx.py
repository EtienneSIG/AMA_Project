"""Build the LearnEU CXO TED-style restitution .pptx.

Reads Restitution01/slides/*.md, recognises a `Render` directive in each
slide's metadata, and chooses one of several visual layouts:

  * hero        -> solid hero background + giant centered title + sub-line
  * stat        -> solid background + huge editable stat number + label
  * quote       -> dark background + pulled quote + attribution
  * image       -> solid background + title + 3-5 bullets overlaid
  * persona     -> 3 rounded-rectangle persona tiles with avatar circles
  * cards       -> classic dark navy cards (original layout, kept for appendices)

Every visual element is a native PowerPoint shape (rectangle, oval,
text box) — no raster images — so the user can edit colour, size,
copy and layout directly in PowerPoint.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "Subject" / "Azure Master Architect_Prezo_Template_v01.pptx"
SLIDES_DIR = ROOT / "Restitution01" / "slides"
ASSETS = ROOT / "Restitution01" / "assets"
OUT = ROOT / "Restitution01" / "build" / "LearnEU-CXO-Restitution.pptx"

# Brand palette (sampled from template chrome)
ORANGE = RGBColor(0xF2, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xC8, 0xD4, 0xE6)
CARD_FILL = RGBColor(0x12, 0x2A, 0x4A)
CARD_LINE = RGBColor(0x2A, 0x4A, 0x78)
NAVY = RGBColor(0x12, 0x2A, 0x4A)
TEAL = RGBColor(0x0E, 0x5A, 0x5A)
PINK = RGBColor(0xB2, 0x3A, 0x6F)
GREEN = RGBColor(0x1F, 0x6B, 0x3A)
INK = RGBColor(0x0F, 0x1B, 0x2D)

SLIDE_W = 12_192_000
SLIDE_H = 6_858_000

TITLE_SLIDE_IDX = 1
CLOSING_SLIDE_IDX = 21  # "Ask of the board"
CLOSING_CAPTION = "What we ask of the board"

# -- Native-shape themes: each background "image" name maps to a solid
# fill + a couple of decorative orbs. The slide markdown still references
# names like "bg-hero-navy.png" so we don't have to rewrite every slide.
BG_THEMES = {
    "bg-hero-navy.png":   {"fill": NAVY,   "orbs": [(ORANGE, 0.35), (TEAL, 0.20)]},
    "bg-hero-orange.png": {"fill": ORANGE, "orbs": [(NAVY,   0.35), (WHITE, 0.10)]},
    "bg-hero-teal.png":   {"fill": TEAL,   "orbs": [(ORANGE, 0.30), (WHITE, 0.10)]},
    "bg-hero-pink.png":   {"fill": PINK,   "orbs": [(ORANGE, 0.30), (WHITE, 0.10)]},
    "bg-hero-green.png":  {"fill": GREEN,  "orbs": [(ORANGE, 0.30), (WHITE, 0.10)]},
    "bg-dark-grid.png":   {"fill": INK,    "orbs": [(NAVY,   0.25)]},
}

# -- Native stat tiles: each "stat-*.png" reference becomes a coloured
# background + a giant editable number + an editable label.
STAT_THEMES = {
    "stat-gap.png":       {"bg": NAVY,   "num": "40%",     "label": "outcome gap, today"},
    "stat-admin.png":     {"bg": ORANGE, "num": "−45%",    "label": "teacher admin time"},
    "stat-12mo.png":      {"bg": TEAL,   "num": "12 → 6",  "label": "months → weeks per market"},
    "stat-learners.png":  {"bg": NAVY,   "num": "4.1M",    "label": "EU learners in scope"},
    "stat-reduction.png": {"bg": ORANGE, "num": "−26pp",   "label": "outcome gap closed"},
    "stat-value.png":     {"bg": GREEN,  "num": "€55M",    "label": "annual run-rate value"},
    "stat-6wk.png":       {"bg": TEAL,   "num": "6 wk",    "label": "localisation per market"},
    "stat-gdpr.png":      {"bg": NAVY,   "num": "100%",    "label": "GDPR Art. 8 compliance"},
}

# -- Native persona tiles
PERSONA_THEMES = {
    "persona-learner.png": {"tint": RGBColor(0x1B, 0x3C, 0x70), "initial": "L",
                             "name": "Lucas", "tag": "12 · Berlin"},
    "persona-teacher.png": {"tint": RGBColor(0x8F, 0x52, 0x10), "initial": "K",
                             "name": "Mr Klein", "tag": "Year-7 maths"},
    "persona-parent.png":  {"tint": RGBColor(0x1F, 0x6B, 0x3A), "initial": "S",
                             "name": "Sophie", "tag": "Lucas' mother"},
}


# ----------------------------- parsing ---------------------------------


def _meta(md: str, name: str) -> str:
    rx = re.compile(rf"^-\s+\*\*{re.escape(name)}.*?:\*\*\s*(.+)$", re.MULTILINE)
    m = rx.search(md)
    return m.group(1).strip() if m else ""


def _section(md: str, name: str) -> str:
    m = re.search(
        rf"^##\s+{re.escape(name)}[^\n]*\n((?:(?!^##\s).+\n?)*)", md, re.MULTILINE
    )
    return m.group(1).strip() if m else ""


def _strip_md(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    return s.strip()


def parse_slide(md: str) -> dict:
    out: dict = {
        "title_line": "",
        "headline": _meta(md, "Headline"),
        "subheadline": _meta(md, "Sub-headline"),
        "layout_hint": _meta(md, "Layout (template)"),
        "render": _meta(md, "Render") or "cards",
        "image": _meta(md, "Image"),
        "images": _meta(md, "Images"),
        "quote": _section(md, "Quote"),
        "attribution": _meta(md, "Attribution"),
        "cxo_focus": _meta(md, "CXO focus"),
        "sources": _meta(md, "Source refs"),
        "visual": _section(md, "Visual"),
        "demo_cue": _section(md, "Demo cue"),
        "notes": _section(md, "Speaker notes"),
        "columns": [],
    }
    m = re.search(r"^#\s+Slide\s+\d+\s*\xb7\s*[^\xb7]+\xb7\s*(.+)$", md, re.MULTILINE)
    if m:
        out["title_line"] = m.group(1).strip()

    for sec in re.finditer(
        r"^##\s+Body bullets(?:\s*\(([^)]*)\))?\s*\n((?:(?!^##\s).+\n?)*)",
        md,
        re.MULTILINE,
    ):
        caption_raw = (sec.group(1) or "").strip()
        caption = ""
        if caption_raw:
            parts = re.split(r"\s*[\u2014\u2013-]\s*", caption_raw, maxsplit=1)
            caption = parts[1].strip().capitalize() if len(parts) == 2 else caption_raw
        bullets: list[str] = []
        for line in sec.group(2).splitlines():
            line = line.strip()
            if line.startswith("- "):
                bullets.append(_strip_md(line[2:]))
        out["columns"].append({"caption": caption, "bullets": bullets})
    return out


# ----------------------------- drawing ---------------------------------


def _set_text(tf, lines: list[tuple[str, int, int, RGBColor, bool]]) -> None:
    tf.word_wrap = True
    tf.margin_left = Emu(120_000)
    tf.margin_right = Emu(120_000)
    tf.margin_top = Emu(80_000)
    tf.margin_bottom = Emu(80_000)
    if not lines:
        tf.text = ""
        return
    first = True
    for text, level, size, color, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_subhead(slide, text: str) -> None:
    if not text:
        return
    box = slide.shapes.add_textbox(
        Emu(640_000), Emu(950_000), Emu(SLIDE_W - 1_280_000), Emu(360_000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _strip_md(text)
    for run in p.runs:
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = SOFT


def _add_card(slide, x, y, w, h, caption: str, bullets: list[str]) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.color.rgb = CARD_LINE
    card.line.width = Emu(12_700)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(220_000)
    tf.margin_right = Emu(220_000)
    tf.margin_top = Emu(180_000)
    tf.margin_bottom = Emu(180_000)
    n = len(bullets)
    body_pt = 18 if n <= 5 else 16 if n <= 7 else 14
    cap_pt = 16
    lines: list[tuple[str, int, int, RGBColor, bool]] = []
    if caption:
        lines.append((caption.upper(), 0, cap_pt, ORANGE, True))
    for b in bullets:
        lines.append(("\u2022 " + b, 0, body_pt, WHITE, False))
    _set_text(tf, lines)


def render_content(slide, columns: list[dict]) -> None:
    body_top = 1_400_000
    body_h = SLIDE_H - body_top - 360_000
    margin = 600_000
    if not columns:
        return
    if len(columns) == 1:
        _add_card(
            slide,
            Emu(margin),
            Emu(body_top),
            Emu(SLIDE_W - 2 * margin),
            Emu(body_h),
            columns[0]["caption"],
            columns[0]["bullets"],
        )
        return
    cols = columns[:3]
    gap = 240_000
    total_w = SLIDE_W - 2 * margin - gap * (len(cols) - 1)
    card_w = total_w // len(cols)
    for i, c in enumerate(cols):
        x = margin + i * (card_w + gap)
        _add_card(
            slide,
            Emu(x),
            Emu(body_top),
            Emu(card_w),
            Emu(body_h),
            c["caption"],
            c["bullets"],
        )


def set_title(slide, text: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = text or ""
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = WHITE


def set_notes(slide, spec: dict) -> None:
    parts: list[str] = []
    if spec["notes"]:
        parts.append(spec["notes"])
    meta = []
    if spec["cxo_focus"]:
        meta.append(f"CXO focus: {spec['cxo_focus']}")
    if spec["visual"]:
        meta.append(f"Visual: {spec['visual']}")
    if spec["demo_cue"]:
        meta.append(f"Demo cue: {spec['demo_cue']}")
    if spec["sources"]:
        meta.append(f"Source refs: {spec['sources']}")
    if meta:
        parts.append("---\n" + "\n".join(meta))
    if not parts:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = "\n\n".join(parts)


def layout_by_name(prs: Presentation, *names: str):
    by = {l.name: l for l in prs.slide_layouts}
    for n in names:
        if n in by:
            return by[n]
    return prs.slide_layouts[0]


def _wipe_other_placeholders(slide) -> None:
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 0:
            continue
        if ph.has_text_frame:
            ph.text_frame.text = ""


def build_title_slide(prs: Presentation, spec: dict):
    layout = layout_by_name(prs, "Title slide")
    s = prs.slides.add_slide(layout)
    set_title(s, spec["headline"])
    _wipe_other_placeholders(s)
    box = s.shapes.add_textbox(
        Emu(640_000), Emu(2_900_000), Emu(SLIDE_W - 1_280_000), Emu(1_400_000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    lines: list[tuple[str, int, int, RGBColor, bool]] = []
    if spec["subheadline"]:
        lines.append((_strip_md(spec["subheadline"]), 0, 22, SOFT, False))
    if spec["columns"]:
        for b in spec["columns"][0]["bullets"]:
            lines.append((b, 0, 16, WHITE, False))
    _set_text(tf, lines)
    return s


def build_closing_slide(prs: Presentation, spec: dict):
    layout = layout_by_name(prs, "1_Closing logo slide", "2_Closing logo slide")
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        if ph.has_text_frame:
            ph.text_frame.text = ""
    box = s.shapes.add_textbox(
        Emu(640_000), Emu(900_000), Emu(SLIDE_W - 1_280_000), Emu(900_000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec["headline"] or ""
    for run in p.runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE
    sub = s.shapes.add_textbox(
        Emu(640_000), Emu(1_900_000), Emu(SLIDE_W - 1_280_000), Emu(500_000)
    )
    sp = sub.text_frame.paragraphs[0]
    sp.text = _strip_md(spec["subheadline"] or "")
    for run in sp.runs:
        run.font.size = Pt(20)
        run.font.italic = True
        run.font.color.rgb = SOFT
    if spec["columns"]:
        body_top = 2_600_000
        body_h = SLIDE_H - body_top - 700_000
        _add_card(
            s,
            Emu(900_000),
            Emu(body_top),
            Emu(SLIDE_W - 1_800_000),
            Emu(body_h),
            CLOSING_CAPTION,
            spec["columns"][0]["bullets"],
        )
    return s


def build_content_slide(prs: Presentation, spec: dict):
    layout = layout_by_name(prs, "Title Only")
    s = prs.slides.add_slide(layout)
    _wipe_other_placeholders(s)
    set_title(s, spec["headline"])
    add_subhead(s, spec["subheadline"])
    render_content(s, spec["columns"])
    return s


# ----------------------------- TED-style layouts -----------------------


def _blank_slide(prs: Presentation):
    """Pick a blank-ish layout and clear it."""
    layout = layout_by_name(prs, "Blank", "Title Only", "Title slide")
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        if ph.has_text_frame:
            ph.text_frame.text = ""
    return s


def _add_shape(slide, shape_type, x, y, w, h, fill, line=None, line_w_pt=None):
    """Add a native PowerPoint auto-shape with solid fill (no picture)."""
    shp = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w_pt is not None:
            shp.line.width = Pt(line_w_pt)
    if shp.has_text_frame:
        shp.text_frame.text = ""
    return shp


def _add_bg(slide, image_name: str) -> None:
    """Draw the slide background as native shapes (rect + decorative orbs).

    `image_name` is kept for backward compatibility with slide markdown
    that still references e.g. `bg-hero-navy.png` — we just look up a
    theme and draw shapes."""
    theme = BG_THEMES.get((image_name or "").strip(), BG_THEMES["bg-hero-navy.png"])
    # Background plate
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["fill"])
    # Decorative orbs (large, soft) — fully editable circles
    orbs = theme.get("orbs", [])
    spots = [
        (SLIDE_W - 2_400_000, -1_200_000, 4_400_000),
        (-1_400_000, SLIDE_H - 2_600_000, 4_000_000),
        (SLIDE_W // 2 + 1_500_000, SLIDE_H - 1_800_000, 2_400_000),
    ]
    for i, (color, _alpha) in enumerate(orbs[:3]):
        x, y, d = spots[i]
        _add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, color)


def _add_textbox(slide, x, y, w, h, text, size, color, bold=False, italic=False,
                 align_center=True):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def build_hero_slide(prs: Presentation, spec: dict):
    """Giant title centred on a full-bleed native background."""
    s = _blank_slide(prs)
    _add_bg(s, spec["image"] or "bg-hero-navy.png")
    title = _strip_md(spec["headline"] or "")
    sub = _strip_md(spec["subheadline"] or "")
    _add_textbox(s, 600_000, 2_300_000, SLIDE_W - 1_200_000, 2_100_000,
                 title, 72, WHITE, bold=True)
    if sub:
        _add_textbox(s, 600_000, 4_400_000, SLIDE_W - 1_200_000, 700_000,
                     sub, 26, SOFT, italic=True)
    return s


def build_stat_slide(prs: Presentation, spec: dict):
    """Solid background + giant editable stat number + label.

    The big number and label are real text boxes — change the wording
    or font directly in PowerPoint."""
    s = _blank_slide(prs)
    img = (spec["image"] or "stat-gap.png").strip()
    theme = STAT_THEMES.get(img, STAT_THEMES["stat-gap.png"])
    # Background plate
    _add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["bg"])
    # Subtle decorative orb in the top-right corner (editable circle)
    _add_shape(s, MSO_SHAPE.OVAL,
               SLIDE_W - 1_400_000, -1_600_000, 3_600_000, 3_600_000, WHITE)
    # Headline ribbon
    if spec["headline"]:
        _add_textbox(s, 600_000, 360_000, SLIDE_W - 1_200_000, 540_000,
                     _strip_md(spec["headline"]), 30, WHITE, bold=True)
    # Big number
    _add_textbox(s, 600_000, 1_700_000, SLIDE_W - 1_200_000, 2_800_000,
                 theme["num"], 240, WHITE, bold=True)
    # Big label
    _add_textbox(s, 600_000, 4_900_000, SLIDE_W - 1_200_000, 700_000,
                 theme["label"], 32, WHITE, italic=True)
    if spec["subheadline"]:
        _add_textbox(s, 600_000, 5_900_000, SLIDE_W - 1_200_000, 540_000,
                     _strip_md(spec["subheadline"]), 22, SOFT, italic=True)
    return s


def build_quote_slide(prs: Presentation, spec: dict):
    """Pull quote on a solid native background."""
    s = _blank_slide(prs)
    _add_bg(s, spec["image"] or "bg-hero-navy.png")
    _add_textbox(s, 800_000, 800_000, 1_400_000, 1_400_000,
                 "\u201C", 220, ORANGE, bold=True, align_center=False)
    quote = _strip_md((spec["quote"] or spec["headline"] or "").strip())
    _add_textbox(s, 1_200_000, 1_900_000, SLIDE_W - 2_400_000, 3_100_000,
                 quote, 44, WHITE, bold=False)
    attr = _strip_md(spec["attribution"] or "")
    if attr:
        _add_textbox(s, 600_000, 5_500_000, SLIDE_W - 1_200_000, 500_000,
                     "\u2014 " + attr, 22, SOFT, italic=True)
    return s


def build_image_slide(prs: Presentation, spec: dict):
    """Solid native background + title strip + up to 5 short bullets."""
    s = _blank_slide(prs)
    _add_bg(s, spec["image"] or "bg-dark-grid.png")
    _add_textbox(s, 600_000, 600_000, SLIDE_W - 1_200_000, 800_000,
                 _strip_md(spec["headline"] or ""), 44, WHITE, bold=True,
                 align_center=False)
    if spec["subheadline"]:
        _add_textbox(s, 600_000, 1_500_000, SLIDE_W - 1_200_000, 540_000,
                     _strip_md(spec["subheadline"]), 22, SOFT, italic=True,
                     align_center=False)
    if spec["columns"]:
        bullets = spec["columns"][0]["bullets"][:5]
        box = slide_textbox(s, 800_000, 2_400_000,
                            SLIDE_W - 1_600_000, SLIDE_H - 2_800_000)
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = "\u2022  " + b
            for run in p.runs:
                run.font.size = Pt(26)
                run.font.color.rgb = WHITE
            p.space_after = Pt(14)
    return s


def slide_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))


def build_persona_slide(prs: Presentation, spec: dict):
    """Three native persona tiles: rounded rectangle + avatar circle + initial."""
    s = _blank_slide(prs)
    _add_bg(s, spec["image"] or "bg-dark-grid.png")
    _add_textbox(s, 600_000, 360_000, SLIDE_W - 1_200_000, 540_000,
                 _strip_md(spec["headline"] or ""), 36, WHITE, bold=True)
    if spec["subheadline"]:
        _add_textbox(s, 600_000, 920_000, SLIDE_W - 1_200_000, 480_000,
                     _strip_md(spec["subheadline"]), 22, SOFT, italic=True)
    imgs = [n.strip() for n in (spec["images"] or "").split(",") if n.strip()]
    if len(imgs) < 3:
        imgs = ["persona-learner.png", "persona-teacher.png", "persona-parent.png"]
    cols = spec["columns"][:3] if spec["columns"] else []
    margin = 400_000
    gap = 280_000
    tile_w = (SLIDE_W - 2 * margin - gap * 2) // 3
    tile_h = 3_400_000
    y = 1_700_000
    for i in range(3):
        x = margin + i * (tile_w + gap)
        theme = PERSONA_THEMES.get(imgs[i], PERSONA_THEMES["persona-learner.png"])
        # Tile background (rounded rectangle)
        _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                   x, y, tile_w, tile_h, theme["tint"])
        # Avatar circle
        avatar_d = 1_500_000
        ax = x + (tile_w - avatar_d) // 2
        ay = y + 300_000
        _add_shape(s, MSO_SHAPE.OVAL, ax, ay, avatar_d, avatar_d, WHITE)
        # Initial letter inside the avatar
        _add_textbox(s, ax, ay + 150_000, avatar_d, avatar_d - 150_000,
                     theme["initial"], 96, theme["tint"], bold=True)
        # Name + tag under the avatar
        _add_textbox(s, x, ay + avatar_d + 200_000, tile_w, 500_000,
                     theme["name"], 28, WHITE, bold=True)
        _add_textbox(s, x, ay + avatar_d + 700_000, tile_w, 400_000,
                     theme["tag"], 16, SOFT, italic=True)
        # Bullets caption below the tile
        if i < len(cols):
            bullets = cols[i]["bullets"][:3]
            box = slide_textbox(s, x, y + tile_h + 100_000, tile_w, 1_400_000)
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for b in bullets:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = b
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.color.rgb = WHITE
    return s


def build_statgrid_slide(prs: Presentation, spec: dict):
    """2x2 grid of huge editable stat tiles.

    Reads up to 4 bullets in the form `<number> | <label>` from the
    slide's columns (left col first, then right col). Each tile is a
    rounded rectangle + a giant number + a label, all native shapes.
    Falls back to plain bullets if the `|` separator is missing.
    """
    s = _blank_slide(prs)
    _add_bg(s, spec["image"] or "bg-hero-navy.png")
    _add_textbox(s, 600_000, 280_000, SLIDE_W - 1_200_000, 540_000,
                 _strip_md(spec["headline"] or ""), 32, WHITE, bold=True)
    if spec["subheadline"]:
        _add_textbox(s, 600_000, 820_000, SLIDE_W - 1_200_000, 480_000,
                     _strip_md(spec["subheadline"]), 20, SOFT, italic=True)

    # Flatten up to 4 bullets across columns
    bullets: list[str] = []
    for c in spec["columns"]:
        bullets.extend(c["bullets"])
    bullets = bullets[:4]

    palette = [ORANGE, TEAL, NAVY, GREEN]
    margin = 600_000
    gap = 280_000
    grid_top = 1_500_000
    grid_h = SLIDE_H - grid_top - 600_000
    tile_w = (SLIDE_W - 2 * margin - gap) // 2
    tile_h = (grid_h - gap) // 2

    for i, b in enumerate(bullets):
        row, col = divmod(i, 2)
        x = margin + col * (tile_w + gap)
        y = grid_top + row * (tile_h + gap)
        if "|" in b:
            number, label = b.split("|", 1)
        elif ":" in b:
            number, label = b.split(":", 1)
        else:
            number, label = b, ""
        number = _strip_md(number).strip()
        label = _strip_md(label).strip()
        _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                   x, y, tile_w, tile_h, palette[i % len(palette)])
        _add_textbox(s, x, y + 200_000, tile_w, tile_h // 2,
                     number, 96, WHITE, bold=True)
        if label:
            _add_textbox(s, x + 300_000, y + tile_h // 2 + 400_000,
                         tile_w - 600_000, tile_h // 2 - 500_000,
                         label, 20, WHITE, italic=True)
    return s


# --------------------------------- main --------------------------------


RENDERERS = {
    "hero": build_hero_slide,
    "stat": build_stat_slide,
    "statgrid": build_statgrid_slide,
    "quote": build_quote_slide,
    "image": build_image_slide,
    "persona": build_persona_slide,
    "cards": build_content_slide,
}


def main() -> int:
    try:
        sys.stdout.reconfigure(encoding="utf-8")  # type: ignore[attr-defined]
    except Exception:
        pass
    if not TEMPLATE.exists():
        print(f"ERROR: template missing: {TEMPLATE}", file=sys.stderr)
        return 2
    prs = Presentation(str(TEMPLATE))

    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    for sldId in list(sldIdLst):
        rId = sldId.attrib[
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        ]
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    md_files = sorted(SLIDES_DIR.glob("slide-*.md"))
    if not md_files:
        print(f"ERROR: no slide files in {SLIDES_DIR}", file=sys.stderr)
        return 2

    for path in md_files:
        m = re.match(r"slide-(\d+)-", path.name)
        if not m:
            continue
        idx = int(m.group(1))
        spec = parse_slide(path.read_text(encoding="utf-8"))
        render = (spec.get("render") or "cards").lower()
        if idx == TITLE_SLIDE_IDX and render == "cards":
            # Default behaviour for the title slide if no Render directive
            slide = build_title_slide(prs, spec)
            ltag = "TitleSlide"
        elif idx == CLOSING_SLIDE_IDX and render == "cards":
            slide = build_closing_slide(prs, spec)
            ltag = "Closing"
        else:
            renderer = RENDERERS.get(render, build_content_slide)
            slide = renderer(prs, spec)
            ltag = render
        set_notes(slide, spec)
        ncols = len(spec["columns"])
        nb = sum(len(c["bullets"]) for c in spec["columns"])
        print(
            f"slide {idx:02d} \xb7 {ltag:10s} \xb7 cols={ncols} bullets={nb} \xb7 "
            f"{(spec['headline'] or spec['title_line'])[:55]}"
        )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(
        f"\nSaved: {OUT}  ({OUT.stat().st_size:,} bytes, {len(md_files)} slides)"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
